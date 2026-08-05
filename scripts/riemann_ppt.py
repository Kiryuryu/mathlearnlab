#!/usr/bin/env python3
"""黎曼和课件 PPT — 含 GIF 动图"""
import os, io, numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.animation import FuncAnimation, PillowWriter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUT_DIR = os.path.dirname(os.path.abspath(__file__))
FIG_DIR = os.path.join(OUT_DIR, '_ppt_figs')
os.makedirs(FIG_DIR, exist_ok=True)

plt.rcParams['font.sans-serif'] = ['Arial Unicode MS', 'SimHei', 'DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False

BG = '#1a2530'
ACCENT = '#4a6a8a'
ACCENT2 = '#d06868'
PALETTE = ['#4a6a8a', '#d06868', '#3d6b4f', '#c4956a', '#7a5aaa', '#5a9a9a']

def f(x): return x**2

# ══════════════════════════════════════════
#  GIF 1: 矩形逐个添加
# ══════════════════════════════════════════
def make_gif_add_rectangles(path):
    fig, ax = plt.subplots(figsize=(7, 4.5))
    x = np.linspace(0, 3, 300)
    ax.plot(x, f(x), color=ACCENT, lw=2.5, zorder=5)
    n_max = 12
    xs_all = np.linspace(0, 3, n_max + 1)
    dx = 3 / n_max
    bars_container = [[]]

    def update(frame):
        for b in bars_container[0]: b.remove()
        bars_container[0] = []
        for i in range(frame + 1):
            xm = xs_all[i] + dx / 2
            h = f(xm)
            bar = ax.bar(xs_all[i] + dx / 2, h, width=dx * 0.85,
                         color=PALETTE[i % len(PALETTE)], alpha=0.55,
                         edgecolor='white', linewidth=0.6, zorder=3)
            bars_container[0].extend(bar)
        approx = sum(f(xs_all[i] + dx / 2) * dx for i in range(frame + 1))
        ax.set_title(f'添加第 {frame+1}/{n_max} 个矩形  累计 ≈ {approx:.3f}',
                     fontsize=13, fontweight='bold', color='white')
        fig.set_facecolor(BG); ax.set_facecolor(BG)
        ax.tick_params(colors='white'); ax.xaxis.label.set_color('white')
        ax.yaxis.label.set_color('white')
        for spine in ax.spines.values(): spine.set_color('#445566')
        return bars_container[0]

    ax.set_xlim(-0.2, 3.3); ax.set_ylim(-0.3, 10.5)
    ax.set_xlabel('$x$', fontsize=12); ax.set_ylabel('$f(x)$', fontsize=12)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    fig.set_facecolor(BG); ax.set_facecolor(BG)
    ax.tick_params(colors='white')
    for spine in ax.spines.values(): spine.set_color('#445566')

    anim = FuncAnimation(fig, update, frames=n_max, interval=600, repeat=True)
    anim.save(path, writer=PillowWriter(fps=2))
    plt.close(fig)
    print(f'  ✓ {path}')

# ══════════════════════════════════════════
#  GIF 2: n 增大，逼近精确值
# ══════════════════════════════════════════
def make_gif_convergence(path):
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 4.5))
    ns = [1, 2, 3, 4, 6, 8, 12, 16, 24, 32, 48, 64]
    exact = 9.0
    x = np.linspace(0, 3, 300)

    def update(frame):
        for ax in (ax1, ax2): ax.clear()
        n = ns[frame]
        dx = 3 / n
        xs = np.linspace(0, 3, n + 1)

        # 左图：柱状图
        ax1.plot(x, f(x), color=ACCENT, lw=2.5, zorder=5)
        for i in range(n):
            xm = xs[i] + dx / 2
            h = f(xm)
            ax1.bar(xs[i] + dx / 2, h, width=dx * 0.85,
                    color=PALETTE[i % len(PALETTE)], alpha=0.5,
                    edgecolor='white', linewidth=0.4, zorder=3)
        approx = sum(f(xs[i] + dx / 2) * dx for i in range(n))
        err = abs(approx - exact)
        ax1.set_title(f'n = {n}  近似 = {approx:.4f}', fontsize=12, fontweight='bold', color='white')
        ax1.set_xlim(-0.2, 3.3); ax1.set_ylim(-0.3, 10.5)
        ax1.set_facecolor(BG)

        # 右图：误差收敛曲线
        ns_so_far = ns[:frame + 1]
        approxs = []
        for nn in ns_so_far:
            ddx = 3 / nn
            xss = np.linspace(0, 3, nn + 1)
            approxs.append(sum(f(xss[i] + ddx / 2) * ddx for i in range(nn)))
        ax2.plot(ns_so_far, approxs, 'o-', color=ACCENT, markersize=5, lw=1.5)
        ax2.axhline(y=exact, color=ACCENT2, ls='--', lw=2, label=f'精确值 = {exact}')
        ax2.set_xlabel('n', fontsize=12, color='white')
        ax2.set_ylabel('近似值', fontsize=12, color='white')
        ax2.set_title(f'误差 = {err:.4f}', fontsize=12, fontweight='bold', color='white')
        ax2.legend(fontsize=10, facecolor='#223344', edgecolor='#445566',
                   labelcolor='white')
        ax2.set_facecolor(BG)

        for ax in (ax1, ax2):
            ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
            ax.tick_params(colors='white')
            for spine in ax.spines.values(): spine.set_color('#445566')
        fig.set_facecolor(BG)
        fig.suptitle('n 越大，近似越精确', fontsize=14, fontweight='bold', color='white', y=1.02)

    anim = FuncAnimation(fig, update, frames=len(ns), interval=800, repeat=True)
    anim.save(path, writer=PillowWriter(fps=1.5))
    plt.close(fig)
    print(f'  ✓ {path}')

# ══════════════════════════════════════════
#  GIF 3: 左/右/中点 对比切换
# ══════════════════════════════════════════
def make_gif_types(path):
    fig, axes = plt.subplots(1, 3, figsize=(12, 4))
    kinds = ['left', 'right', 'mid']
    labels = ['左端点', '右端点', '中点']
    n = 8
    x = np.linspace(0, 3, 300)
    xs = np.linspace(0, 3, n + 1)
    dx = 3 / n

    def update(frame):
        for ax in axes: ax.clear()
        kind = kinds[frame]
        for idx, (ax, k, lab) in enumerate(zip(axes, kinds, labels)):
            ax.plot(x, f(x), color=ACCENT, lw=2.5, zorder=5)
            for i in range(n):
                if k == 'left':   xm = xs[i]
                elif k == 'right': xm = xs[i] + dx
                else:              xm = xs[i] + dx / 2
                h = f(xm)
                bar_color = PALETTE[i % 5] if k == kind else '#55667744'
                alpha = 0.55 if k == kind else 0.2
                ax.bar(xs[i] + dx / 2, h, width=dx * 0.85,
                       color=bar_color, alpha=alpha,
                       edgecolor='white', linewidth=0.5, zorder=3)
            approx = sum((f(xs[i]) if k=='left' else f(xs[i]+dx) if k=='right' else f(xs[i]+dx/2)) * dx for i in range(n))
            fontw = 'bold' if k == kind else 'normal'
            ax.set_title(f'{lab}  ≈ {approx:.3f}', fontsize=12, fontweight=fontw, color='white')
            ax.set_xlim(-0.2, 3.3); ax.set_ylim(-0.3, 10.5)
            ax.set_facecolor(BG)
            ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
            ax.tick_params(colors='white')
            for spine in ax.spines.values(): spine.set_color('#445566')
        fig.set_facecolor(BG)

    anim = FuncAnimation(fig, update, frames=3, interval=1200, repeat=True)
    anim.save(path, writer=PillowWriter(fps=1))
    plt.close(fig)
    print(f'  ✓ {path}')

# ══════════════════════════════════════════
#  静态图
# ══════════════════════════════════════════
def fig_integral(path):
    fig, ax = plt.subplots(figsize=(7, 4.5))
    x = np.linspace(0, 3, 300)
    ax.plot(x, f(x), color=ACCENT, lw=3, zorder=5)
    ax.fill_between(x, f(x), alpha=0.2, color=ACCENT, zorder=2)
    ax.text(1.5, 2.5, r'$\int_0^3 x^2\,dx = 9$', fontsize=20, ha='center',
            color='white', fontweight='bold',
            bbox=dict(boxstyle='round,pad=0.4', facecolor='#223344', edgecolor=ACCENT, alpha=0.95))
    ax.set_xlim(-0.3, 3.3); ax.set_ylim(-0.5, 10.5)
    ax.set_xlabel('$x$', fontsize=13); ax.set_ylabel('$f(x)$', fontsize=13)
    ax.set_title('定积分 = 黎曼和的极限', fontsize=14, fontweight='bold', color='white')
    ax.set_facecolor(BG); fig.set_facecolor(BG)
    ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
    ax.tick_params(colors='white')
    for spine in ax.spines.values(): spine.set_color('#445566')
    fig.tight_layout()
    fig.savefig(path, dpi=180, bbox_inches='tight', facecolor=BG)
    plt.close(fig)
    print(f'  ✓ {path}')

# ══════════════════════════════════════════
#  PPT 构建
# ══════════════════════════════════════════
def add_bg(slide, color_hex=BG.replace('#','')):
    bg = slide.background; fill = bg.fill; fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color_hex)

def add_title(slide, text, top=0.4):
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(8.8), Inches(0.7))
    tf = txBox.text_frame; p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(32); p.font.color.rgb = RGBColor.from_string('FFFFFF')
    p.font.bold = True; p.alignment = PP_ALIGN.LEFT

def add_bar(slide, top=0.35):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(top), Inches(1.2), Inches(0.06))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor.from_string(ACCENT.replace('#',''))
    shape.line.fill.background()

def add_text(slide, text, top=1.2, left=0.6, w=8.8, size=20, color='CCCCCC'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(5))
    tf = txBox.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(size)
        p.font.color.rgb = RGBColor.from_string(color); p.space_after = Pt(6)

def build():
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(7.5)

    # ── 1. 封面 ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title(s, '黎曼和', 2.2)
    add_bar(s, 3.15)
    add_text(s, 'Riemann Sum\n从矩形近似到定积分', 3.5, size=24, color='AABBCC')
    add_text(s, '数学博物馆 · 数学之美展区', 6.5, size=14, color='667788')

    # ── 2. 什么是黎曼和 + GIF ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title(s, '什么是黎曼和？'); add_bar(s)
    add_text(s,
        '核心思想：用矩形面积之和近似曲线下面积\n\n'
        '把 [a, b] 分成 n 个小区间\n'
        '每个区间上画一个矩形\n'
        '所有矩形面积相加 ≈ 曲线下面积',
        1.3, size=20)
    gif1 = os.path.join(FIG_DIR, 'add_rectangles.gif')
    s.shapes.add_picture(gif1, Inches(4.5), Inches(1.2), Inches(5.2), Inches(3.3))
    # 公式
    txBox = s.shapes.add_textbox(Inches(4.6), Inches(4.7), Inches(5), Inches(1.5))
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = 'S = Σ f(xᵢ*) · Δxᵢ'
    p.font.size = Pt(24); p.font.color.rgb = RGBColor.from_string('FFFFFF')
    p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = 'Δx = (b−a)/n'
    p2.font.size = Pt(14); p2.font.color.rgb = RGBColor.from_string('8899AA')
    p2.alignment = PP_ALIGN.CENTER

    # ── 3. 三种取点 + GIF ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title(s, '三种取点方式'); add_bar(s)
    add_text(s,
        '左端点 → 低估递增函数\n'
        '右端点 → 高估递增函数\n'
        '中  点 → 通常更精确',
        1.2, size=20)
    gif3 = os.path.join(FIG_DIR, 'types.gif')
    s.shapes.add_picture(gif3, Inches(0.2), Inches(3.0), Inches(9.6), Inches(4.2))

    # ── 4. n 增大收敛 + GIF ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title(s, 'n 越大，近似越精确'); add_bar(s)
    gif2 = os.path.join(FIG_DIR, 'convergence.gif')
    s.shapes.add_picture(gif2, Inches(0.1), Inches(1.1), Inches(9.8), Inches(4.5))
    add_text(s,
        '∫₀³ x² dx = 9\n'
        'n=1 → 6.75  |  n=4 → 8.44  |  n=16 → 8.88  |  n=64 → 8.99',
        5.8, left=1, w=8, size=18, color='FFFFFF')

    # ── 5. 定积分 ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title(s, '从黎曼和到定积分'); add_bar(s)
    fi = os.path.join(FIG_DIR, 'integral.png')
    s.shapes.add_picture(fi, Inches(0.3), Inches(1.1), Inches(9.4), Inches(4.2))
    add_text(s,
        '牛顿-莱布尼茨公式：\n'
        '∫ₐᵇ f(x) dx = F(b) − F(a)    其中 F\'(x) = f(x)',
        5.5, left=2, w=6, size=22, color='FFFFFF')

    # ── 6. 动手试试 ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title(s, '动手试试'); add_bar(s)
    add_text(s,
        '例：∫₀² (x + 1) dx\n\n'
        '分 [0,2] 为 n 份，Δx = 2/n，取右端点 xᵢ = 2i/n\n\n'
        'Sₙ = Σ (2i/n + 1)·(2/n)\n'
        '   = 2(n+1)/n + 2\n\n'
        'lim Sₙ = 2 + 2 = 4  ✓\n'
        '验证：[x²/2 + x]₀² = 2 + 2 = 4',
        1.2, size=20)

    # ── 7. 总结 ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title(s, '总结'); add_bar(s)
    add_text(s,
        '✦ 黎曼和 = 用矩形近似曲线下面积\n\n'
        '✦ 三种取点：左端点 / 右端点 / 中点\n\n'
        '✦ 矩形越多 (n→∞)，近似越精确\n\n'
        '✦ 黎曼和的极限 = 定积分\n\n'
        '✦ 定积分是微积分的核心概念之一',
        1.3, size=24)

    # ── 8. 谢谢 ──
    s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s)
    add_title(s, '谢谢', 3); add_bar(s, 3.95)
    add_text(s, '数学博物馆 · mathlearnlab.cn', 4.3, size=18, color='667788')

    out = os.path.join(OUT_DIR, '黎曼和课件.pptx')
    prs.save(out)
    print(f'\n✅ PPT: {out}')

if __name__ == '__main__':
    print('生成 GIF 动图...')
    make_gif_add_rectangles(os.path.join(FIG_DIR, 'add_rectangles.gif'))
    make_gif_convergence(os.path.join(FIG_DIR, 'convergence.gif'))
    make_gif_types(os.path.join(FIG_DIR, 'types.gif'))
    print('生成静态图...')
    fig_integral(os.path.join(FIG_DIR, 'integral.png'))
    print('\n构建 PPT...')
    build()
